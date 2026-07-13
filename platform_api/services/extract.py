"""Document text extraction for supported file types."""

from __future__ import annotations

from pathlib import Path


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".docx":
        from docx import Document

        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            lines.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
        return "\n".join(lines)
    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"# Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)
    raise ValueError(f"unsupported file type: {suffix}")
